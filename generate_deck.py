"""
generate_deck.py
================
Generates a 17-slide PowerPoint presentation for the Capital Markets AI Safety Platform.

Usage:
    python generate_deck.py [--out FILEPATH]

Requirements:
    pip install python-pptx

Output:
    CapitalMarkets_AI_Safety_Platform.pptx  (or the path specified by --out)

Deck structure
--------------
Section 1 – Account Team (6 slides)
  1.  Title slide
  2.  Business Problem & Risk Landscape
  3.  Platform Overview (capability wheel)
  4.  Solution Architecture
  5.  Human-in-the-Loop Compliance Controls
  6.  Business Value & ROI

Section 2 – Specialist Team (4 slides)
  7.  Azure AI Content Safety Services Matrix
  8.  Content Filters & Guardrail Engineering
  9.  Azure AI Foundry Control Plane
  10. Compliance Pipeline – Parallel Execution

Section 3 – Solution Engineer (7 slides)
  11. Technical Architecture Deep-Dive
  12. Backend Services Reference
  13. Frontend Component Map
  14. Demo Agent Fleet (10 archetypes)
  15. Configuration & Environment Variables
  16. Deployment Guide
  17. Screenshots & Live Demo Environments
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx is required.  Run: pip install python-pptx")
    sys.exit(1)


# ---------------------------------------------------------------------------
# Brand palette
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x0D, 0x1B, 0x2E)
CYAN   = RGBColor(0x00, 0xB4, 0xD8)
GOLD   = RGBColor(0xFF, 0xC4, 0x00)
GREEN  = RGBColor(0x00, 0xE5, 0x8A)
RED    = RGBColor(0xFF, 0x5E, 0x57)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x8E, 0xA0, 0xB8)
LIGHT  = RGBColor(0xE8, 0xED, 0xF5)
SLATE  = RGBColor(0x1A, 0x2A, 0x40)
TEAL   = RGBColor(0x06, 0xB6, 0xD4)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)


# ---------------------------------------------------------------------------
# Slide dimensions (16:9)
# ---------------------------------------------------------------------------
W = Inches(13.33)
H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width_pt=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                 font_name="Segoe UI", wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_slide(prs, layout_idx=6):
    """Add a blank slide."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # Remove all placeholder shapes
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)
    return slide


def dark_background(slide):
    add_rect(slide, 0, 0, W, H, NAVY)


def section_badge(slide, text, color=CYAN):
    bw = Inches(2.4)
    bh = Inches(0.28)
    add_rect(slide, Inches(0.5), Inches(0.28), bw, bh, color)
    add_text_box(slide, text.upper(),
                 Inches(0.5), Inches(0.27), bw, bh,
                 font_size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


def slide_title(slide, title, subtitle=None, title_color=WHITE, sub_color=GREY):
    add_text_box(slide, title,
                 Inches(0.5), Inches(0.65), Inches(12.0), Inches(0.7),
                 font_size=28, bold=True, color=title_color)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.5), Inches(1.3), Inches(11.0), Inches(0.4),
                     font_size=14, bold=False, color=sub_color)


def divider(slide, top_inch=1.85):
    add_rect(slide, Inches(0.5), Inches(top_inch), Inches(12.3), Inches(0.02), CYAN)


def footer(slide, page_num, total=17):
    add_text_box(slide, f"Capital Markets AI Safety Platform   |   {page_num}/{total}",
                 Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3),
                 font_size=8, color=GREY, align=PP_ALIGN.RIGHT)


def bullet_block(slide, title, bullets, left, top, width, height,
                 title_color=CYAN, bullet_color=LIGHT, title_size=11, bullet_size=9.5,
                 accent_color=None):
    accent = accent_color or CYAN
    add_rect(slide, left, top, Inches(0.04), height, accent)
    add_text_box(slide, title,
                 left + Inches(0.12), top, width - Inches(0.12), Inches(0.3),
                 font_size=title_size, bold=True, color=title_color)
    bullet_top = top + Inches(0.3)
    for b in bullets:
        add_text_box(slide, f"  {b}",
                     left + Inches(0.12), bullet_top,
                     width - Inches(0.12), Inches(0.22),
                     font_size=bullet_size, color=bullet_color)
        bullet_top += Inches(0.21)


def kpi_card(slide, label, value, left, top, width, height,
             card_color=SLATE, value_color=CYAN, label_color=GREY):
    add_rect(slide, left, top, width, height, card_color)
    add_text_box(slide, value,
                 left, top + Inches(0.12), width, Inches(0.38),
                 font_size=26, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_text_box(slide, label,
                 left, top + Inches(0.50), width, Inches(0.24),
                 font_size=9, color=label_color, align=PP_ALIGN.CENTER)


def service_pill(slide, name, detail, left, top, accent=CYAN):
    pw = Inches(3.8)
    ph = Inches(0.52)
    add_rect(slide, left, top, pw, ph, SLATE)
    add_rect(slide, left, top, Inches(0.03), ph, accent)
    add_text_box(slide, name,
                 left + Inches(0.1), top + Inches(0.04), pw - Inches(0.1), Inches(0.22),
                 font_size=9.5, bold=True, color=WHITE)
    add_text_box(slide, detail,
                 left + Inches(0.1), top + Inches(0.27), pw - Inches(0.1), Inches(0.2),
                 font_size=8, color=GREY)


# ---------------------------------------------------------------------------
# Slide 1: Title
# ---------------------------------------------------------------------------

def slide_01_title(prs):
    slide = add_slide(prs)
    dark_background(slide)

    # Left accent bar
    add_rect(slide, 0, 0, Inches(0.06), H, CYAN)

    # Logo area top-right
    add_text_box(slide, "Microsoft Azure",
                 Inches(10.5), Inches(0.35), Inches(2.5), Inches(0.35),
                 font_size=10, color=GREY, align=PP_ALIGN.RIGHT)

    # Main title
    add_text_box(slide,
                 "Capital Markets\nAI Safety Platform",
                 Inches(0.8), Inches(1.6), Inches(9.0), Inches(1.8),
                 font_size=40, bold=True, color=WHITE)

    # Tagline
    add_text_box(slide,
                 "End-to-End AI Governance for Capital Markets",
                 Inches(0.8), Inches(3.4), Inches(9.0), Inches(0.5),
                 font_size=18, color=CYAN)

    # Sub-tagline
    add_text_box(slide,
                 "Detect  |  Block  |  Audit  |  Remediate",
                 Inches(0.8), Inches(3.95), Inches(9.0), Inches(0.4),
                 font_size=13, color=GREY)

    # Right panel
    add_rect(slide, Inches(9.8), Inches(1.4), Inches(3.3), Inches(4.8), SLATE)
    add_text_box(slide, "DESIGNED FOR",
                 Inches(10.0), Inches(1.55), Inches(3.0), Inches(0.25),
                 font_size=8, bold=True, color=CYAN, align=PP_ALIGN.LEFT)
    personas = [
        "Head of AI Governance",
        "Chief AI Risk Officer",
        "Compliance Technology Lead",
        "Capital Markets CTO / CDO",
        "Financial Regulatory Teams",
    ]
    for i, p in enumerate(personas):
        add_text_box(slide, f"  {p}",
                     Inches(10.0), Inches(1.85) + Inches(0.38) * i,
                     Inches(3.0), Inches(0.32),
                     font_size=9.5, color=LIGHT)

    # Bottom bar
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), SLATE)
    add_text_box(slide,
                 "Azure AI Content Safety  |  Azure AI Foundry  |  Azure OpenAI  |  Azure AI Language",
                 Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.25),
                 font_size=8, color=GREY, align=PP_ALIGN.CENTER)

    footer(slide, 1)
    return slide


# ---------------------------------------------------------------------------
# Slide 2: Business Problem
# ---------------------------------------------------------------------------

def slide_02_business_problem(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Account Team — Business Context", CYAN)
    slide_title(slide, "The Capital Markets AI Risk Gap",
                "Four critical failures that expose firms to regulatory and reputational risk")
    divider(slide)
    footer(slide, 2)

    risks = [
        ("Insider Trading via AI",
         RED,
         ["AI assistants draft communications referencing MNPI",
          "Distributed to clients before compliance review",
          "MiFID II Art. 14 / SEC Rule 10b-5 exposure",
          "SEC inquiry opened; remediation cost: $2M+"]),
        ("Jailbreak & Override Attacks",
         AMBER,
         ["Adversaries bypass AI safety controls",
          "Trading AIs instructed to execute unauthorised orders",
          "Traditional keyword blocklists fail semantic attacks",
          "No audit trail for blocked attempts"]),
        ("AI Hallucination in Research",
         GOLD,
         ["Research notes cite metrics not in source data",
          "Fabricated financials distributed to institutional clients",
          "FINRA Rule 2210 supervisory failures",
          "Class action liability for misleading communications"]),
        ("AI Fleet Blind Spot",
         PURPLE,
         ["Guardrail status unknown until quarterly audit",
          "30+ AI agents running without centralised oversight",
          "Compliance posture: reactive, not real-time",
          "No single view of regulatory control coverage"]),
    ]

    for i, (title, color, bullets, *_) in enumerate(risks):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * Inches(6.2)
        top = Inches(2.1) + row * Inches(2.3)
        bw = Inches(5.8)
        bh = Inches(2.1)
        add_rect(slide, left, top, bw, bh, SLATE)
        add_rect(slide, left, top, Inches(0.05), bh, color)
        add_text_box(slide, title,
                     left + Inches(0.15), top + Inches(0.1), bw, Inches(0.3),
                     font_size=11, bold=True, color=WHITE)
        for j, b in enumerate(bullets):
            add_text_box(slide, f"  {b}",
                         left + Inches(0.15), top + Inches(0.42) + Inches(0.38) * j,
                         bw - Inches(0.2), Inches(0.3),
                         font_size=9, color=LIGHT)

    return slide


# ---------------------------------------------------------------------------
# Slide 3: Platform overview
# ---------------------------------------------------------------------------

def slide_03_platform_overview(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Account Team — Platform Overview", CYAN)
    slide_title(slide, "One Platform. Nine Safety Layers.",
                "Every AI-generated communication passes through a unified governance stack")
    divider(slide)
    footer(slide, 3)

    services = [
        ("Text Analysis",       "Hate / Violence / SelfHarm / Sexual — severity 0-6",     CYAN),
        ("Image Analysis",      "Visual content moderation across all 4 harm categories",  TEAL),
        ("Prompt Shields",      "Jailbreak & XPIA (indirect injection) detection",         CYAN),
        ("Groundedness",        "Hallucination detection against source documents",         GREEN),
        ("Protected Material",  "Copyright and licensed IP detection",                      TEAL),
        ("Custom Categories",   "Insider Trading, Market Manipulation, Front Running",      AMBER),
        ("Blocklist Manager",   "Sanctions, restricted securities, prohibited terms",       GOLD),
        ("Task Adherence",      "AI agent scope enforcement via gpt-4o evaluation",         PURPLE),
        ("PII Detection",       "12+ entity types — SSN, CreditCard, IBAN, and more",      RED),
    ]

    for i, (name, detail, color) in enumerate(services):
        col = i % 3
        row = i // 3
        left = Inches(0.5) + col * Inches(4.1)
        top = Inches(2.1) + row * Inches(0.66)
        service_pill(slide, name, detail, left, top, accent=color)

    # Compliance pipeline callout (bottom)
    add_rect(slide, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.85), SLATE)
    add_text_box(slide, "COMPLIANCE PIPELINE",
                 Inches(0.65), Inches(6.15), Inches(3.0), Inches(0.25),
                 font_size=9, bold=True, color=CYAN)
    add_text_box(slide,
                 "Orchestrates all 9 services in parallel via asyncio.gather — "
                 "single weighted risk score (0-100) in under 800ms.  "
                 "PASS < 30  |  REVIEW 30-69  |  BLOCK >= 70",
                 Inches(0.65), Inches(6.42), Inches(12.0), Inches(0.45),
                 font_size=9, color=LIGHT)

    return slide


# ---------------------------------------------------------------------------
# Slide 4: Solution architecture
# ---------------------------------------------------------------------------

def slide_04_architecture(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Account Team — Solution Architecture", CYAN)
    slide_title(slide, "Architecture Overview",
                "React + FastAPI + Azure AI — fully stateless, demo-mode by default")
    divider(slide)
    footer(slide, 4)

    layers = [
        ("FRONTEND", "React 18 + Vite 5.4, port 5173",
         "27+ routes  |  Recharts  |  lucide-react  |  inline CSS", CYAN, Inches(0.5)),
        ("BACKEND API", "FastAPI 0.115.4 + uvicorn, port 8000",
         "5 routers  |  asyncio.gather  |  DEMO_MODE fallback", TEAL, Inches(3.2)),
        ("AZURE AI SERVICES", "Azure AI Content Safety  |  Azure OpenAI  |  Azure AI Language",
         "SDK 1.0.0  |  REST 2024-02-15-preview  |  gpt-4o  |  LM 2023-04-01", GREEN, Inches(5.9)),
        ("AZURE AI FOUNDRY", "Control plane: azure-ai-projects 1.0.0b10  |  mgmt SDK 13.5.0",
         "Guardrails  |  Model deployments  |  Agent fleet  |  Security alerts", PURPLE, Inches(8.6)),
        ("AUTH & IDENTITY", "azure-identity 1.19.0  |  ClientSecretCredential",
         "DEMO_MODE bypasses all credential requirements", AMBER, Inches(11.3)),
    ]

    for label, line1, line2, color, left in layers:
        w = Inches(2.4)
        h = Inches(2.2)
        top = Inches(2.2)
        add_rect(slide, left, top, w, h, SLATE)
        add_rect(slide, left, top, w, Inches(0.04), color)
        add_text_box(slide, label,
                     left + Inches(0.1), top + Inches(0.1), w, Inches(0.22),
                     font_size=7.5, bold=True, color=color)
        add_text_box(slide, line1,
                     left + Inches(0.1), top + Inches(0.36), w - Inches(0.1), Inches(0.55),
                     font_size=8.5, bold=True, color=WHITE)
        add_text_box(slide, line2,
                     left + Inches(0.1), top + Inches(0.95), w - Inches(0.1), Inches(1.0),
                     font_size=8, color=GREY)

    # Data flow arrows (text representation)
    add_text_box(slide,
                 "Browser  ->  React SPA  ->  FastAPI  ->  Azure AI (parallel)  ->  Response",
                 Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.3),
                 font_size=9, color=GREY, align=PP_ALIGN.CENTER)

    # Demo mode callout
    add_rect(slide, Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.65), _rgb(16, 20, 35))
    add_rect(slide, Inches(0.5), Inches(5.15), Inches(0.04), Inches(0.65), GREEN)
    add_text_box(slide, "DEMO MODE (default: DEMO_MODE=true)",
                 Inches(0.65), Inches(5.2), Inches(5.0), Inches(0.25),
                 font_size=9, bold=True, color=GREEN)
    add_text_box(slide,
                 "All 9 services return synthetic capital markets data from mock_data.py.  "
                 "No Azure credentials required.  Production mode activated by setting DEMO_MODE=false and providing .env credentials.",
                 Inches(0.65), Inches(5.45), Inches(12.0), Inches(0.3),
                 font_size=8.5, color=LIGHT)

    return slide


# ---------------------------------------------------------------------------
# Slide 5: HITL controls
# ---------------------------------------------------------------------------

def slide_05_hitl(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Account Team — Compliance Controls", CYAN)
    slide_title(slide, "Human-in-the-Loop Compliance Gates",
                "Every high-risk decision requires explicit human authorisation — by design")
    divider(slide)
    footer(slide, 5)

    gates = [
        ("REVIEW Gate",
         "Compliance Pipeline",
         "Risk score 30-69 halts auto-routing. "
         "Compliance officer must explicitly Approve or Reject. "
         "Decision logged with timestamp and user ID.",
         AMBER,
         "MiFID II Art. 16 (supervisory controls)"),
        ("Remediation Confirmation",
         "Foundry Control Plane > Model Deployments",
         "One-click guardrail remediation requires confirmation dialog before applying.  "
         "Prevents accidental policy changes on production deployments.",
         RED,
         "FINRA Rule 3110 (supervisory systems)"),
        ("System Prompt Lock",
         "Agent Registry (Settings)",
         "All system prompts and guardrail configurations are read-only in the UI.  "
         "Changes require engineering review — prevents unauthorised scope expansion.",
         PURPLE,
         "SEC Rule 17a-4 (audit trail)"),
        ("Guardrail Required Flag",
         "Foundry Control Plane > Compliance Policies",
         "Policies flag any deployment with disabled controls.  "
         "Violation Detected status triggers mandatory remediation workflow.",
         GREEN,
         "Basel III Pillar 2 (operational risk)"),
    ]

    for i, (gate, location, desc, color, regulation) in enumerate(gates):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * Inches(6.2)
        top = Inches(2.15) + row * Inches(2.4)
        bw = Inches(5.8)
        bh = Inches(2.2)
        add_rect(slide, left, top, bw, bh, SLATE)
        add_rect(slide, left, top, Inches(0.05), bh, color)
        add_text_box(slide, gate,
                     left + Inches(0.15), top + Inches(0.08), bw, Inches(0.28),
                     font_size=11.5, bold=True, color=WHITE)
        add_text_box(slide, f"Location: {location}",
                     left + Inches(0.15), top + Inches(0.38), bw - Inches(0.2), Inches(0.22),
                     font_size=8.5, color=color)
        add_text_box(slide, desc,
                     left + Inches(0.15), top + Inches(0.64), bw - Inches(0.2), Inches(1.0),
                     font_size=9, color=LIGHT)
        add_text_box(slide, regulation,
                     left + Inches(0.15), top + Inches(1.88), bw - Inches(0.2), Inches(0.22),
                     font_size=8, color=GREY)

    return slide


# ---------------------------------------------------------------------------
# Slide 6: Business Value
# ---------------------------------------------------------------------------

def slide_06_value(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Account Team — Business Value", GOLD)
    slide_title(slide, "Business Value & ROI",
                "From 48-hour reactive compliance review to sub-800ms proactive enforcement")
    divider(slide)
    footer(slide, 6)

    kpis = [
        ("800 ms", "Pipeline end-to-end latency (9 parallel services)", CYAN),
        ("0 MNPI", "AI-generated MNPI communications reaching clients", GREEN),
        ("100%", "Fleet guardrail coverage — real-time, not quarterly audit", GOLD),
        ("48h -> 0", "Manual review cycle eliminated for clean content", TEAL),
    ]

    for i, (val, lbl, color) in enumerate(kpis):
        kpi_card(slide, lbl, val,
                 Inches(0.5) + Inches(3.15) * i, Inches(2.1),
                 Inches(2.9), Inches(1.1),
                 value_color=color)

    comparisons = [
        ("OLD APPROACH", RED, [
            "48-hour manual review per AI communication",
            "Insider trading language detected post-distribution",
            "Guardrail status: unknown until quarterly audit",
            "Custom AI risk detectors: 6-12 month ML project",
            "Jailbreak attacks reach AI models unfiltered",
        ]),
        ("THIS PLATFORM", GREEN, [
            "Sub-800ms automated pipeline, 9 parallel safety checks",
            "MNPI blocked before distribution with regulatory citation",
            "Real-time fleet dashboard, one-click remediation",
            "Four financial crime categories deployed in minutes",
            "Prompt Shields block adversarial inputs in <200ms",
        ]),
    ]

    for i, (title, color, bullets) in enumerate(comparisons):
        left = Inches(0.5) + i * Inches(6.2)
        top = Inches(3.5)
        bw = Inches(5.8)
        bh = Inches(3.7)
        add_rect(slide, left, top, bw, bh, SLATE)
        add_rect(slide, left, top, bw, Inches(0.04), color)
        add_text_box(slide, title,
                     left + Inches(0.15), top + Inches(0.12), bw, Inches(0.28),
                     font_size=11, bold=True, color=color)
        for j, b in enumerate(bullets):
            add_text_box(slide, f"  {b}",
                         left + Inches(0.15), top + Inches(0.5) + Inches(0.58) * j,
                         bw - Inches(0.2), Inches(0.5),
                         font_size=9.5, color=LIGHT)

    return slide


# ---------------------------------------------------------------------------
# Slide 7: Content Safety Services Matrix
# ---------------------------------------------------------------------------

def slide_07_cs_matrix(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Specialist Team — Content Safety", TEAL)
    slide_title(slide, "Azure AI Content Safety — Services Matrix",
                "All 9 services with API endpoint, SDK version, and capital markets use case")
    divider(slide)
    footer(slide, 7)

    headers = ["Service", "API / SDK", "Use Case", "Latency (demo)"]
    col_widths = [Inches(2.5), Inches(3.2), Inches(4.5), Inches(1.6)]
    col_lefts = [Inches(0.5), Inches(3.05), Inches(6.3), Inches(10.85)]
    row_h = Inches(0.44)
    header_top = Inches(2.1)

    # Header row
    add_rect(slide, Inches(0.5), header_top, Inches(12.3), row_h, SLATE)
    for j, (h, w, l) in enumerate(zip(headers, col_widths, col_lefts)):
        add_text_box(slide, h, l, header_top + Inches(0.1), w, Inches(0.28),
                     font_size=9, bold=True, color=CYAN)

    rows = [
        ("Text Analysis",       "azure-ai-contentsafety 1.0.0",        "Harm categories: Hate, Violence, Sexual, SelfHarm",
         "~140 ms"),
        ("Image Analysis",      "azure-ai-contentsafety 1.0.0",        "Visual content moderation for uploaded documents",
         "~310 ms"),
        ("Prompt Shields",      "REST 2024-02-15-preview",             "Jailbreak + XPIA indirect injection detection",
         "~200 ms"),
        ("Groundedness",        "REST 2024-02-15-preview",             "Hallucination detection vs. source documents",
         "~900 ms"),
        ("Protected Material",  "REST 2024-02-15-preview",             "Copyright / licensed IP detection in outputs",
         "~280 ms"),
        ("Custom Categories",   "Incidents API 2024-02-15-preview",    "Insider trading, market manipulation, front running",
         "~350 ms"),
        ("Blocklist Manager",   "azure-ai-contentsafety 1.0.0 SDK",    "Sanctions, restricted ISINs, prohibited trading terms",
         "~120 ms"),
        ("Task Adherence",      "Azure OpenAI gpt-4o",                 "Agent scope enforcement — advisory vs. execution roles",
         "~1800 ms"),
        ("PII Detection",       "Azure AI Language REST 2023-04-01",   "SSN, CreditCard, IBAN, Email, PhoneNumber (12+ types)",
         "~220 ms"),
    ]

    for i, (svc, api, use, lat) in enumerate(rows):
        top_i = header_top + row_h * (i + 1)
        row_bg = _rgb(26, 34, 55) if i % 2 == 0 else SLATE
        add_rect(slide, Inches(0.5), top_i, Inches(12.3), row_h, row_bg)
        vals = [svc, api, use, lat]
        colors = [WHITE, GREY, LIGHT, CYAN]
        for j, (v, c, w, l) in enumerate(zip(vals, colors, col_widths, col_lefts)):
            add_text_box(slide, v, l, top_i + Inches(0.1),
                         w, Inches(0.28), font_size=9, color=c)

    return slide


# ---------------------------------------------------------------------------
# Slide 8: Guardrail engineering
# ---------------------------------------------------------------------------

def slide_08_guardrails(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Specialist Team — Content Filters", AMBER)
    slide_title(slide, "Guardrail Engineering — Content Filters",
                "Build, test, compare, and measure guardrails via the Azure AI Foundry data-plane API")
    divider(slide)
    footer(slide, 8)

    bullet_block(
        slide, "Guardrail Manager",
        ["Three default profiles: Permissive / Standard / Strict",
         "CRUD via Foundry data-plane API v2025-05-15-preview",
         "Push new guardrail config in seconds, not sprints",
         "All changes logged for audit trail"],
        Inches(0.5), Inches(2.2), Inches(5.8), Inches(2.0),
        accent_color=AMBER,
    )

    bullet_block(
        slide, "Model Filter Test — 7 Pre-Built Scenarios",
        ["clean_query: baseline compliant capital markets question",
         "jailbreak_trade: DAN-style override to execute $50M short",
         "market_manipulation: coordinated spoofing language",
         "insider_trading: MNPI-laced earnings preview",
         "regulatory_bypass: CTR structuring instructions",
         "indirect_injection: XPIA via document content",
         "threat_content: violence directed at committee members"],
        Inches(6.6), Inches(2.2), Inches(6.0), Inches(2.8),
        accent_color=RED,
    )

    bullet_block(
        slide, "Filter Comparison",
        ["Run same scenario through two guardrail configs simultaneously",
         "Side-by-side ALLOWED vs. BLOCKED verdict comparison",
         "Compliance evidence artefact: before/after guardrail posture"],
        Inches(0.5), Inches(4.5), Inches(5.8), Inches(1.8),
        accent_color=TEAL,
    )

    bullet_block(
        slide, "Filter Analytics",
        ["Block rate chart per category over time (Recharts)",
         "Entity heatmap: which agents generate most block events",
         "In-memory event store fed by all test and pipeline runs",
         "Answers FINRA Rule 3110 supervisory record requirement"],
        Inches(6.6), Inches(5.25), Inches(6.0), Inches(1.8),
        accent_color=PURPLE,
    )

    return slide


# ---------------------------------------------------------------------------
# Slide 9: Foundry control plane
# ---------------------------------------------------------------------------

def slide_09_foundry(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Specialist Team — Foundry Control Plane", PURPLE)
    slide_title(slide, "Azure AI Foundry Control Plane",
                "Real-time governance of the entire AI fleet — agents, deployments, policies, security alerts")
    divider(slide)
    footer(slide, 9)

    pages = [
        ("Fleet Overview",
         CYAN,
         ["Fleet health score (0-100) aggregated across all agents",
          "KPIs: Total Agents / Compliance % / Alerts / Daily Cost",
          "Hourly API volume chart (Recharts)",
          "Activity feed: recent policy violations and remediations"]),
        ("Agent Fleet",
         GREEN,
         ["Platform filter: Foundry / AutoGen / Semantic Kernel / LangChain / Custom",
          "Status filter: All / Active / Warning / Error",
          "Health score, compliance status, token usage per agent",
          "Required by MiFID II Art. 16 and FINRA Rule 3110"]),
        ("Model Deployments",
         AMBER,
         ["Three guardrail indicators per deployment (CF / PS / AM)",
          "Red badge on any deployment with disabled control",
          "Quota usage bar: token consumption vs. limit",
          "One-click Remediate button with confirmation gate"]),
        ("Compliance Policies",
         PURPLE,
         ["Maps deployments to: MiFID II / FINRA / GDPR / SOC2",
          "Controls checklist per policy framework",
          "Violation count and total assets in scope",
          "Compliance posture score aggregated across all policies"]),
    ]

    for i, (page, color, bullets) in enumerate(pages):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * Inches(6.2)
        top = Inches(2.15) + row * Inches(2.5)
        bw = Inches(5.8)
        bh = Inches(2.3)
        add_rect(slide, left, top, bw, bh, SLATE)
        add_rect(slide, left, top, Inches(0.05), bh, color)
        add_text_box(slide, page,
                     left + Inches(0.15), top + Inches(0.08),
                     bw, Inches(0.28), font_size=11.5, bold=True, color=WHITE)
        for j, b in enumerate(bullets):
            add_text_box(slide, f"  {b}",
                         left + Inches(0.15), top + Inches(0.45) + Inches(0.44) * j,
                         bw - Inches(0.2), Inches(0.38), font_size=9, color=LIGHT)

    return slide


# ---------------------------------------------------------------------------
# Slide 10: Compliance pipeline
# ---------------------------------------------------------------------------

def slide_10_pipeline(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Specialist Team — Pipeline Architecture", GREEN)
    slide_title(slide, "Compliance Pipeline — Parallel Execution",
                "asyncio.gather fan-out / fan-in — 9 services, single weighted verdict")
    divider(slide)
    footer(slide, 10)

    # Weight table
    weights = [
        ("Prompt Shields",    20, CYAN),
        ("Market Manipulation", 20, AMBER),
        ("Insider Trading",   20, RED),
        ("Front Running",     15, _rgb(249, 115, 22)),
        ("Text Analysis",     15, TEAL),
        ("Protected Material", 10, PURPLE),
        ("Groundedness",      10, GREEN),
    ]

    add_text_box(slide, "Service Weights (sum = 110 — normalised to 0-100)",
                 Inches(0.5), Inches(2.15), Inches(5.5), Inches(0.28),
                 font_size=9.5, bold=True, color=CYAN)

    for i, (svc, weight, color) in enumerate(weights):
        top = Inches(2.5) + Inches(0.53) * i
        add_rect(slide, Inches(0.5), top, Inches(5.5), Inches(0.42), SLATE)
        add_text_box(slide, svc, Inches(0.65), top + Inches(0.1), Inches(2.5), Inches(0.25),
                     font_size=9, color=LIGHT)
        bar_w = Inches(weight / 25)
        add_rect(slide, Inches(3.3), top + Inches(0.14), bar_w, Inches(0.15), color)
        add_text_box(slide, str(weight), Inches(3.3) + bar_w + Inches(0.05), top + Inches(0.1),
                     Inches(0.3), Inches(0.25), font_size=8.5, color=color)

    # Threshold callouts
    thresholds = [
        ("PASS",   "Score < 30",  GREEN,  Inches(6.5), Inches(2.15)),
        ("REVIEW", "Score 30-69", AMBER,  Inches(6.5), Inches(2.9)),
        ("BLOCK",  "Score >= 70", RED,    Inches(6.5), Inches(3.65)),
    ]
    add_text_box(slide, "Verdict Thresholds",
                 Inches(6.5), Inches(2.15) - Inches(0.32), Inches(5.5), Inches(0.28),
                 font_size=9.5, bold=True, color=CYAN)
    for verdict, rule, color, left, top in thresholds:
        add_rect(slide, left, top, Inches(5.5), Inches(0.6), SLATE)
        add_rect(slide, left, top, Inches(0.05), Inches(0.6), color)
        add_text_box(slide, verdict, left + Inches(0.15), top + Inches(0.1),
                     Inches(1.2), Inches(0.28), font_size=12, bold=True, color=color)
        add_text_box(slide, rule, left + Inches(1.4), top + Inches(0.15),
                     Inches(3.5), Inches(0.28), font_size=10, color=GREY)

    # Pipeline facts
    bullet_block(
        slide, "Key Pipeline Facts",
        ["asyncio.gather — all services run concurrently, not sequentially",
         "Wall-clock latency: <800ms for all 9 services (demo mode)",
         "BLOCK verdict halts routing; REVIEW triggers HITL gate",
         "Each service result stored with latency for audit log",
         "Optional: Groundedness requires grounding source document",
         "Compliance pipeline endpoint: POST /api/compliance/pipeline"],
        Inches(6.5), Inches(4.55), Inches(6.3), Inches(2.5),
        accent_color=CYAN,
    )

    return slide


# ---------------------------------------------------------------------------
# Slide 11: Technical architecture deep-dive
# ---------------------------------------------------------------------------

def slide_11_tech_arch(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Solution Engineer — Technical Architecture", PURPLE)
    slide_title(slide, "Technical Architecture Deep-Dive",
                "Service topology, async patterns, and in-memory state management")
    divider(slide)
    footer(slide, 11)

    bullet_block(
        slide, "Backend (FastAPI / Python)",
        ["backend/main.py: app entry, lifespan provisions blocklist + incidents + guardrails",
         "5 APIRouters: /api/content-safety, /api/compliance, /api/foundry, /api/content-filters, /api/demo",
         "asyncio.WindowsSelectorEventLoopPolicy() set for Windows compat",
         "All service functions are async def with httpx.AsyncClient",
         "TTL cache (120s) on Foundry management data via functools.lru_cache",
         "DEMO_MODE=true: returns data from mock_data.py, no Azure calls"],
        Inches(0.5), Inches(2.15), Inches(5.9), Inches(3.0),
        accent_color=CYAN,
    )

    bullet_block(
        slide, "Frontend (React / Vite)",
        ["src/App.jsx: all 27+ routes, react-router-dom 6 BrowserRouter",
         "Layout: Header + Sidebar (4 navigation sections) + main content area",
         "No Tailwind — pure inline CSS using CSS custom properties (--bg-base, --accent-blue, etc.)",
         "Recharts 2.13 for all analytical charts (BarChart, LineChart, RadialBar)",
         "lucide-react 0.462 for all icons",
         "Vite 5.4 dev server with /api proxy to FastAPI on port 8000"],
        Inches(6.7), Inches(2.15), Inches(6.0), Inches(3.0),
        accent_color=PURPLE,
    )

    bullet_block(
        slide, "State & Data Flow",
        ["All API state: useState + useEffect (no global store / Redux)",
         "Content filter events stored in-memory list in content_filters.py service",
         "Foundry data TTL-cached in Python process (not persisted across restarts)",
         "No database — fully stateless, designed for demo repeatability"],
        Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.5),
        accent_color=GREEN,
    )

    return slide


# ---------------------------------------------------------------------------
# Slide 12: Backend services reference
# ---------------------------------------------------------------------------

def slide_12_backend(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Solution Engineer — Backend Services", TEAL)
    slide_title(slide, "Backend Services Reference",
                "14 service modules in backend/services/")
    divider(slide)
    footer(slide, 12)

    services = [
        ("text_analysis.py",      "AnalyzeTextAction SDK — 4 harm categories, severity 0-6",          CYAN),
        ("image_analysis.py",     "AnalyzeImageOptions SDK — base64 or URL input",                    TEAL),
        ("prompt_shields.py",     "REST POST /shieldPrompt — userPromptAttackResults + docs",          CYAN),
        ("groundedness.py",       "REST POST /groundedness:detect — grounding sources list",           GREEN),
        ("protected_material.py", "REST POST /text:detectProtectedMaterial — copyright detection",    TEAL),
        ("custom_categories.py",  "Incidents API — create/hydrate/analyse — 3 financial categories",  AMBER),
        ("blocklist.py",          "BlocklistClient SDK — CRUD operations + analysis scan",             GOLD),
        ("task_adherence.py",     "gpt-4o evaluation: task definition + conversation + tool calls",   PURPLE),
        ("pii_detection.py",      "Azure AI Language REST — recognizeEntities + PII domain",          RED),
        ("content_filters.py",    "Foundry data-plane REST — guardrail CRUD + inference tests",       AMBER),
        ("foundry_mgmt.py",       "azure-ai-projects + mgmt SDK — agents, deployments, policies",     PURPLE),
        ("mock_data.py",          "Synthetic data — 30 agents, 20 deployments, 15 alerts (seeded)",   GREY),
        ("_transport.py",         "AsyncClient factory with timeout and Content-Type defaults",        GREY),
    ]

    for i, (fname, detail, color) in enumerate(services):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * Inches(6.3)
        top = Inches(2.15) + row * Inches(0.48)
        add_rect(slide, left, top, Inches(6.0), Inches(0.42), SLATE)
        add_rect(slide, left, top, Inches(0.03), Inches(0.42), color)
        add_text_box(slide, fname,
                     left + Inches(0.1), top + Inches(0.05),
                     Inches(2.0), Inches(0.22), font_size=8.5, bold=True, color=WHITE,
                     font_name="Consolas")
        add_text_box(slide, detail,
                     left + Inches(2.15), top + Inches(0.08),
                     Inches(3.7), Inches(0.28), font_size=8, color=GREY)

    return slide


# ---------------------------------------------------------------------------
# Slide 13: Frontend component map
# ---------------------------------------------------------------------------

def slide_13_frontend(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Solution Engineer — Frontend Components", CYAN)
    slide_title(slide, "Frontend Component Map",
                "src/components/ and src/pages/ — 40+ React components")
    divider(slide)
    footer(slide, 13)

    sections = [
        ("Layout/",         TEAL,   ["Header.jsx — route-based titles via ROUTE_TITLES",
                                     "Sidebar.jsx — 4 nav sections, icon + label",
                                     "Layout.jsx — flex wrapper"]),
        ("Dashboard/",      CYAN,   ["Dashboard.jsx — KPI cards, Recharts overview"]),
        ("ContentSafety/",  GREEN,  ["TextAnalysis, ImageAnalysis, PromptShields",
                                     "Groundedness, ProtectedMaterial, PIIDetection",
                                     "CustomCategories, Blocklist, TaskAdherence"]),
        ("ContentFilters/", AMBER,  ["GuardrailManager, ModelFilterTest, AgentFilterTest",
                                     "FilterComparison, FilterAnalytics, FilterTestPage",
                                     "JailbreakFilter, PIIFilter, XPIAFilter"]),
        ("FoundryControl/", PURPLE, ["FoundryOverview, AgentFleet, ModelDeployments",
                                     "CompliancePolicies, SecurityAlerts, QuotaManagement, AdminProjects"]),
        ("Common/",         GREY,   ["FeaturePage.jsx — generic feature wrapper"]),
    ]

    for i, (section, color, bullets) in enumerate(sections):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * Inches(6.2)
        top = Inches(2.15) + row * Inches(1.7)
        bw = Inches(5.8)
        bh = Inches(1.55)
        add_rect(slide, left, top, bw, bh, SLATE)
        add_rect(slide, left, top, Inches(0.05), bh, color)
        add_text_box(slide, section,
                     left + Inches(0.15), top + Inches(0.08),
                     bw, Inches(0.28), font_size=10.5, bold=True, color=WHITE,
                     font_name="Consolas")
        for j, b in enumerate(bullets):
            add_text_box(slide, f"  {b}",
                         left + Inches(0.15), top + Inches(0.42) + Inches(0.35) * j,
                         bw - Inches(0.2), Inches(0.3), font_size=9, color=LIGHT)

    return slide


# ---------------------------------------------------------------------------
# Slide 14: Agent fleet
# ---------------------------------------------------------------------------

def slide_14_agents(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Solution Engineer — Demo Agent Fleet", PURPLE)
    slide_title(slide, "Demo Agent Fleet — 10 Archetypes",
                "mock_data.py provides 30 synthetic agents drawn from these 10 role archetypes")
    divider(slide)
    footer(slide, 14)

    agents = [
        ("Trade Research Copilot",    "Azure AI Foundry", "gpt-4o",       "High",     "MiFID II Art. 20, FINRA 2210"),
        ("Risk Assessment Agent",     "Semantic Kernel",  "gpt-4o",       "High",     "Basel III, EMIR, Dodd-Frank"),
        ("Portfolio Optimizer",       "Azure AI Foundry", "o1-preview",   "High",     "MiFID II RTS 28, UCITS"),
        ("Compliance Monitor",        "Azure AI Foundry", "gpt-4o",       "Critical", "FINRA 3110, MiFID II Art. 16"),
        ("Client Advisory Bot",       "LangChain",        "gpt-4o-mini",  "Medium",   "FCA COBS 4, GDPR Art. 5"),
        ("Market Intelligence Agent", "AutoGen",          "gpt-4o",       "Medium",   "EU MAR, SEC Reg FD"),
        ("Regulatory Screening Agent","Azure AI Foundry", "gpt-4o",       "Critical", "FATF Rec. 10, OFAC, AML"),
        ("Fixed Income Analyst",      "Semantic Kernel",  "gpt-4o-mini",  "Medium",   "MiFID II Art. 20, ESMA"),
        ("Equity Research Bot",       "AutoGen",          "gpt-4o",       "High",     "MiFID II Art. 20, FINRA 2241"),
        ("Derivatives Pricing Agent", "Custom",           "Phi-4",        "High",     "EMIR, Dodd-Frank Title VII"),
    ]

    headers = ["Agent", "Platform", "Model", "Risk", "Regulatory Controls"]
    col_lefts = [Inches(0.5), Inches(3.2), Inches(4.95), Inches(6.5), Inches(7.5)]
    col_widths = [Inches(2.65), Inches(1.65), Inches(1.45), Inches(0.9), Inches(5.3)]
    row_h = Inches(0.42)
    header_top = Inches(2.15)

    add_rect(slide, Inches(0.5), header_top, Inches(12.3), row_h, SLATE)
    for j, (h, w, l) in enumerate(zip(headers, col_widths, col_lefts)):
        add_text_box(slide, h, l, header_top + Inches(0.1), w, Inches(0.25),
                     font_size=9, bold=True, color=CYAN)

    risk_colors = {"Critical": RED, "High": AMBER, "Medium": TEAL}
    for i, (name, platform, model, risk, regs) in enumerate(agents):
        top_i = header_top + row_h * (i + 1)
        row_bg = _rgb(26, 34, 55) if i % 2 == 0 else SLATE
        add_rect(slide, Inches(0.5), top_i, Inches(12.3), row_h, row_bg)
        rc = risk_colors.get(risk, GREY)
        vals = [name, platform, model, risk, regs]
        colors = [WHITE, GREY, GREY, rc, LIGHT]
        bolds = [True, False, False, True, False]
        for j, (v, c, b, w, l) in enumerate(zip(vals, colors, bolds, col_widths, col_lefts)):
            add_text_box(slide, v, l, top_i + Inches(0.1),
                         w, Inches(0.25), font_size=8.5, color=c, bold=b)

    return slide


# ---------------------------------------------------------------------------
# Slide 15: Configuration reference
# ---------------------------------------------------------------------------

def slide_15_config(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Solution Engineer — Configuration", AMBER)
    slide_title(slide, "Configuration & Environment Variables",
                "backend/config.py — Pydantic Settings, all values from .env")
    divider(slide)
    footer(slide, 15)

    env_vars = [
        ("DEMO_MODE",                "true",  "true: all services use mock_data.py, no Azure calls", GREEN),
        ("AZURE_CS_ENDPOINT",        "",      "Azure AI Content Safety resource endpoint URL", CYAN),
        ("AZURE_CS_KEY",             "",      "Content Safety API key", CYAN),
        ("AZURE_CS_ENDPOINT_REAL",   "",      "Override endpoint (bypasses fallback chain)", TEAL),
        ("AZURE_OPENAI_ENDPOINT",    "",      "Azure OpenAI resource endpoint URL", PURPLE),
        ("AZURE_OPENAI_KEY",         "",      "Azure OpenAI API key", PURPLE),
        ("AZURE_LANGUAGE_ENDPOINT",  "",      "Azure AI Language resource endpoint", TEAL),
        ("AZURE_LANGUAGE_KEY",       "",      "Azure AI Language API key", TEAL),
        ("AZURE_TENANT_ID",          "",      "Azure AD tenant ID (for ClientSecretCredential)", AMBER),
        ("AZURE_CLIENT_ID",          "",      "Service principal client ID", AMBER),
        ("AZURE_CLIENT_SECRET",      "",      "Service principal client secret", AMBER),
        ("AZURE_SUBSCRIPTION_ID",    "",      "Azure subscription ID (for mgmt SDK)", GOLD),
        ("AZURE_RESOURCE_GROUP",     "",      "Resource group containing Foundry / CS resources", GOLD),
        ("AZURE_FOUNDRY_PROJECT",    "",      "Azure AI Foundry project resource name", PURPLE),
    ]

    for i, (key, default, desc, color) in enumerate(env_vars):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * Inches(6.3)
        top = Inches(2.15) + row * Inches(0.48)
        add_rect(slide, left, top, Inches(6.0), Inches(0.42), SLATE)
        add_rect(slide, left, top, Inches(0.03), Inches(0.42), color)
        add_text_box(slide, key,
                     left + Inches(0.1), top + Inches(0.05),
                     Inches(2.6), Inches(0.22), font_size=8, bold=True, color=color,
                     font_name="Consolas")
        add_text_box(slide, desc,
                     left + Inches(2.75), top + Inches(0.08),
                     Inches(3.1), Inches(0.28), font_size=8, color=GREY)

    return slide


# ---------------------------------------------------------------------------
# Slide 16: Deployment guide
# ---------------------------------------------------------------------------

def slide_16_deployment(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Solution Engineer — Deployment", GREEN)
    slide_title(slide, "Deployment Guide",
                "7 steps from clone to live demo — under 10 minutes")
    divider(slide)
    footer(slide, 16)

    steps = [
        ("1", "Clone the repository",
         "git clone <repo>  &&  cd contentsafety",
         CYAN),
        ("2", "Install Python dependencies",
         "pip install -r requirements.txt  (Python 3.11+ required)",
         TEAL),
        ("3", "Install frontend dependencies",
         "cd frontend  &&  npm install",
         GREEN),
        ("4", "Configure environment (optional for demo mode)",
         "Copy .env.example to .env — skip for DEMO_MODE=true",
         AMBER),
        ("5", "Start the backend",
         "run_backend.bat  (or: uvicorn backend.main:app --reload --port 8000)",
         CYAN),
        ("6", "Start the frontend",
         "run_frontend.bat  (or: cd frontend && npm run dev)",
         PURPLE),
        ("7", "Open the application",
         "Navigate to http://localhost:5173 in Chrome",
         GREEN),
    ]

    for i, (num, title, cmd, color) in enumerate(steps):
        col = i % 2
        row = i // 2
        if i == 6:
            col = 0
            row = 3
        left = Inches(0.5) + col * Inches(6.2)
        top = Inches(2.15) + row * Inches(1.2)
        bw = Inches(5.8) if i < 6 else Inches(12.3)
        bh = Inches(1.05)
        add_rect(slide, left, top, bw, bh, SLATE)
        add_rect(slide, left, top, Inches(0.05), bh, color)
        # Step number circle
        add_rect(slide, left + Inches(0.15), top + Inches(0.2), Inches(0.35), Inches(0.35), color)
        add_text_box(slide, num,
                     left + Inches(0.15), top + Inches(0.17), Inches(0.35), Inches(0.35),
                     font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, title,
                     left + Inches(0.62), top + Inches(0.08), bw - Inches(0.7), Inches(0.28),
                     font_size=10.5, bold=True, color=WHITE)
        add_text_box(slide, cmd,
                     left + Inches(0.62), top + Inches(0.42), bw - Inches(0.7), Inches(0.45),
                     font_size=8.5, color=GREY, font_name="Consolas")

    return slide


# ---------------------------------------------------------------------------
# Slide 17: Screenshots & live environments
# ---------------------------------------------------------------------------

def slide_17_screenshots(prs):
    slide = add_slide(prs)
    dark_background(slide)
    section_badge(slide, "Solution Engineer — Demo Environments", TEAL)
    slide_title(slide, "Live Demo Environments",
                "Key feature pages to navigate during a live demo")
    divider(slide)
    footer(slide, 17)

    demo_pages = [
        ("Compliance Pipeline", "http://localhost:5173",
         "Navigate: Compliance Pipeline",
         "Run clean baseline then insider trading scenario — show score gauge and BLOCK verdict",
         CYAN),
        ("Prompt Shields", "http://localhost:5173/prompt-shields",
         "Navigate: Content Safety > Prompt Shields",
         "Select DAN mode jailbreak scenario — show user_prompt_attack_detected: true",
         TEAL),
        ("Custom Categories", "http://localhost:5173/custom-categories",
         "Navigate: Content Safety > Custom Categories",
         "Run Market Manipulation scenario — show financial crime Incidents API detection",
         AMBER),
        ("Model Deployments", "http://localhost:5173/foundry/deployments",
         "Navigate: Foundry > Model Deployments",
         "Show red badge on noncompliant deployment — click Remediate with confirmation gate",
         PURPLE),
        ("Guardrail Manager", "http://localhost:5173/content-filters/guardrails",
         "Navigate: Content Filters > Guardrail Manager",
         "Show three default profiles — Permissive / Standard / Strict",
         GREEN),
        ("Agent Registry", "http://localhost:5173/settings",
         "Navigate: Design > Agent Registry",
         "Click Trade Research Copilot — expand system prompt — show guardrail config read-only lock",
         RED),
    ]

    for i, (page, url, nav, action, color) in enumerate(demo_pages):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * Inches(6.2)
        top = Inches(2.15) + row * Inches(1.65)
        bw = Inches(5.8)
        bh = Inches(1.5)
        add_rect(slide, left, top, bw, bh, SLATE)
        add_rect(slide, left, top, Inches(0.05), bh, color)
        add_text_box(slide, page,
                     left + Inches(0.15), top + Inches(0.08), bw, Inches(0.28),
                     font_size=11.5, bold=True, color=WHITE)
        add_text_box(slide, nav,
                     left + Inches(0.15), top + Inches(0.38), bw - Inches(0.2), Inches(0.22),
                     font_size=8.5, color=color, font_name="Consolas")
        add_text_box(slide, action,
                     left + Inches(0.15), top + Inches(0.65), bw - Inches(0.2), Inches(0.72),
                     font_size=9, color=LIGHT)

    return slide


# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------

def build(out_path: str):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    print("Building slides...")
    slide_01_title(prs)
    print("  [1/17] Title")
    slide_02_business_problem(prs)
    print("  [2/17] Business Problem")
    slide_03_platform_overview(prs)
    print("  [3/17] Platform Overview")
    slide_04_architecture(prs)
    print("  [4/17] Architecture")
    slide_05_hitl(prs)
    print("  [5/17] HITL Controls")
    slide_06_value(prs)
    print("  [6/17] Business Value")
    slide_07_cs_matrix(prs)
    print("  [7/17] Content Safety Matrix")
    slide_08_guardrails(prs)
    print("  [8/17] Guardrail Engineering")
    slide_09_foundry(prs)
    print("  [9/17] Foundry Control Plane")
    slide_10_pipeline(prs)
    print("  [10/17] Compliance Pipeline")
    slide_11_tech_arch(prs)
    print("  [11/17] Technical Architecture")
    slide_12_backend(prs)
    print("  [12/17] Backend Services")
    slide_13_frontend(prs)
    print("  [13/17] Frontend Components")
    slide_14_agents(prs)
    print("  [14/17] Agent Fleet")
    slide_15_config(prs)
    print("  [15/17] Configuration")
    slide_16_deployment(prs)
    print("  [16/17] Deployment Guide")
    slide_17_screenshots(prs)
    print("  [17/17] Live Demo Environments")

    prs.save(out_path)
    print(f"\nDeck saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate Capital Markets AI Safety Platform deck")
    parser.add_argument("--out", default="CapitalMarkets_AI_Safety_Platform.pptx",
                        help="Output path for the .pptx file (default: CapitalMarkets_AI_Safety_Platform.pptx)")
    args = parser.parse_args()
    build(args.out)
